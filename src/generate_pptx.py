"""Generate stock_deck.pptx from aggregated.json."""

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── paths ────────────────────────────────────────────────────────────────────
INPUT  = Path(__file__).parent / "outputs" / "aggregated.json"
OUTPUT = Path(__file__).parent / "outputs" / "stock_deck.pptx"

# ── palette ──────────────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT       = RGBColor(0x00, 0xB4, 0xD8)   # cyan
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY   = RGBColor(0xCC, 0xD6, 0xE0)
CARD_BG      = RGBColor(0x17, 0x2A, 0x3E)   # slightly lighter navy
GREEN        = RGBColor(0x06, 0xD6, 0x6A)
ORANGE       = RGBColor(0xFF, 0x6B, 0x35)
RED          = RGBColor(0xFF, 0x2D, 0x55)
YELLOW       = RGBColor(0xFF, 0xD1, 0x66)

RISK_COLORS = {
    "Conservative": GREEN,
    "Moderate":     YELLOW,
    "Aggressive":   ORANGE,
    "Speculative":  RED,
}

SPEED_EMOJI = {
    "Steady Climber": "📈",
    "Fast Mover":     "🚀",
    "Rocket":         "🔥",
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── helpers ──────────────────────────────────────────────────────────────────

def rgb_hex(rgb: RGBColor) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, x, y, w, h, *,
                size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    return txb


def add_rect(slide, x, y, w, h, fill: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.line.fill.background()          # no border
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def accent_bar(slide, y=Inches(0.08)):
    """Thin cyan bar at top of every slide."""
    add_rect(slide, 0, y, SLIDE_W, Inches(0.04), ACCENT)


# ── slide builders ───────────────────────────────────────────────────────────

def slide_title(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, DARK_BG)
    accent_bar(slide, y=Inches(0))

    # big title
    add_textbox(slide, "Stock Performance Analysis",
                Inches(1), Inches(2), Inches(11.33), Inches(1.2),
                size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # subtitle
    add_textbox(slide, "ATL → ATH Journey  ·  10 Tickers  ·  10-Year Daily Bars (Split-Adjusted)",
                Inches(1), Inches(3.2), Inches(11.33), Inches(0.6),
                size=18, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    # cyan divider
    add_rect(slide, Inches(4.5), Inches(3.95), Inches(4.33), Inches(0.04), ACCENT)

    # data source line
    add_textbox(slide, "Source: Yahoo Finance  ·  Not investment advice",
                Inches(1), Inches(4.3), Inches(11.33), Inches(0.4),
                size=12, color=LIGHT_GREY, italic=True, align=PP_ALIGN.CENTER)

    # ticker chips row
    tickers = [d["ticker"] for d in data]
    chip_w  = Inches(1.0)
    total_w = chip_w * len(tickers) + Inches(0.15) * (len(tickers) - 1)
    start_x = (SLIDE_W - total_w) / 2
    for i, t in enumerate(tickers):
        x = start_x + i * (chip_w + Inches(0.15))
        add_rect(slide, x, Inches(5.2), chip_w, Inches(0.45), CARD_BG)
        add_textbox(slide, t,
                    x, Inches(5.2), chip_w, Inches(0.45),
                    size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


def slide_overview_table(prs, data):
    """Summary table: all 10 tickers on one slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    accent_bar(slide)

    add_textbox(slide, "Overview — All Tickers",
                Inches(0.4), Inches(0.2), Inches(12), Inches(0.55),
                size=26, bold=True, color=WHITE)

    headers = ["Ticker", "ATL Price", "ATL Date", "ATH Price", "ATH Date", "Days", "Speed", "Risk"]
    col_w   = [Inches(1.0), Inches(1.2), Inches(1.3), Inches(1.2), Inches(1.3),
               Inches(0.8), Inches(1.6), Inches(1.5)]
    row_h   = Inches(0.52)
    start_x = Inches(0.3)
    start_y = Inches(0.9)

    # header row
    x = start_x
    for i, (hdr, cw) in enumerate(zip(headers, col_w)):
        add_rect(slide, x, start_y, cw, row_h, ACCENT)
        add_textbox(slide, hdr, x + Inches(0.05), start_y + Inches(0.07),
                    cw - Inches(0.1), row_h,
                    size=11, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        x += cw

    # data rows
    for row_i, d in enumerate(data):
        y   = start_y + row_h * (row_i + 1)
        bg  = CARD_BG if row_i % 2 == 0 else DARK_BG
        x   = start_x
        mult = d["ath_price"] / d["atl_price"]
        vals = [
            d["ticker"],
            f"${d['atl_price']:,.2f}",
            d["atl_date"],
            f"${d['ath_price']:,.2f}",
            d["ath_date"],
            str(d["days_between"]),
            d["speed_label"],
            d["risk_label"],
        ]
        for i, (val, cw) in enumerate(zip(vals, col_w)):
            add_rect(slide, x, y, cw, row_h, bg)
            color = RISK_COLORS.get(d["risk_label"], WHITE) if i == 7 else \
                    ACCENT if i == 0 else LIGHT_GREY
            add_textbox(slide, val, x + Inches(0.05), y + Inches(0.08),
                        cw - Inches(0.1), row_h,
                        size=10, bold=(i == 0), color=color, align=PP_ALIGN.CENTER)
            x += cw


def slide_multiplier_chart(prs, data):
    """Horizontal bar chart — price multiplier per ticker."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    accent_bar(slide)

    add_textbox(slide, "Price Multiplier: ATH ÷ ATL",
                Inches(0.4), Inches(0.2), Inches(10), Inches(0.55),
                size=26, bold=True, color=WHITE)
    add_textbox(slide, "(how many times higher was the ATH vs ATL)",
                Inches(0.4), Inches(0.72), Inches(10), Inches(0.35),
                size=13, color=LIGHT_GREY, italic=True)

    sorted_data = sorted(data, key=lambda d: d["ath_price"] / d["atl_price"])
    max_mult    = max(d["ath_price"] / d["atl_price"] for d in sorted_data)
    chart_x     = Inches(1.6)
    chart_right = Inches(12.8)
    chart_w     = chart_right - chart_x
    bar_h       = Inches(0.42)
    gap         = Inches(0.12)
    start_y     = Inches(1.2)

    for i, d in enumerate(sorted_data):
        y    = start_y + i * (bar_h + gap)
        mult = d["ath_price"] / d["atl_price"]
        bar_len = chart_w * (mult / max_mult)
        color   = RISK_COLORS.get(d["risk_label"], ACCENT)

        # label
        add_textbox(slide, d["ticker"],
                    Inches(0.3), y + Inches(0.05), Inches(1.2), bar_h,
                    size=12, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
        # bar
        add_rect(slide, chart_x, y, max(bar_len, Inches(0.15)), bar_h, color)
        # value
        add_textbox(slide, f"{mult:.1f}×",
                    chart_x + bar_len + Inches(0.08), y + Inches(0.05),
                    Inches(1.2), bar_h,
                    size=11, bold=True, color=WHITE)


def slide_days_chart(prs, data):
    """Horizontal bar — days between ATL and ATH."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    accent_bar(slide)

    add_textbox(slide, "Days from ATL to ATH",
                Inches(0.4), Inches(0.2), Inches(10), Inches(0.55),
                size=26, bold=True, color=WHITE)
    add_textbox(slide, "(fewer days = faster journey to all-time high)",
                Inches(0.4), Inches(0.72), Inches(10), Inches(0.35),
                size=13, color=LIGHT_GREY, italic=True)

    sorted_data = sorted(data, key=lambda d: d["days_between"])
    max_days    = max(d["days_between"] for d in sorted_data)
    chart_x     = Inches(1.6)
    chart_w     = Inches(11.2)
    bar_h       = Inches(0.42)
    gap         = Inches(0.12)
    start_y     = Inches(1.2)

    for i, d in enumerate(sorted_data):
        y       = start_y + i * (bar_h + gap)
        bar_len = chart_w * (d["days_between"] / max_days)
        color   = ACCENT

        add_textbox(slide, d["ticker"],
                    Inches(0.3), y + Inches(0.05), Inches(1.2), bar_h,
                    size=12, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
        add_rect(slide, chart_x, y, max(bar_len, Inches(0.15)), bar_h, color)
        add_textbox(slide, f"{d['days_between']:,} days",
                    chart_x + bar_len + Inches(0.08), y + Inches(0.05),
                    Inches(1.5), bar_h,
                    size=11, color=WHITE)


def slide_risk_breakdown(prs, data):
    """Risk label distribution with big stat cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    accent_bar(slide)

    add_textbox(slide, "Risk Profile Distribution",
                Inches(0.4), Inches(0.2), Inches(12), Inches(0.55),
                size=26, bold=True, color=WHITE)

    from collections import Counter
    counts = Counter(d["risk_label"] for d in data)
    labels = ["Conservative", "Moderate", "Aggressive", "Speculative"]

    card_w = Inches(2.8)
    card_h = Inches(2.2)
    gap    = Inches(0.4)
    total_row_w = card_w * 4 + gap * 3
    start_x = (SLIDE_W - total_row_w) / 2
    y = Inches(1.1)

    for i, label in enumerate(labels):
        x     = start_x + i * (card_w + gap)
        color = RISK_COLORS[label]
        count = counts.get(label, 0)
        tickers_in = [d["ticker"] for d in data if d["risk_label"] == label]

        add_rect(slide, x, y, card_w, card_h, CARD_BG)
        # color top stripe
        add_rect(slide, x, y, card_w, Inches(0.12), color)
        # count
        add_textbox(slide, str(count),
                    x, y + Inches(0.2), card_w, Inches(0.9),
                    size=52, bold=True, color=color, align=PP_ALIGN.CENTER)
        # label
        add_textbox(slide, label,
                    x, y + Inches(1.1), card_w, Inches(0.45),
                    size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # tickers
        add_textbox(slide, "  ".join(tickers_in) if tickers_in else "—",
                    x + Inches(0.1), y + Inches(1.6), card_w - Inches(0.2), Inches(0.5),
                    size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    # bottom: per-ticker risk table
    y2 = Inches(3.6)
    add_textbox(slide, "Ticker Risk Labels",
                Inches(0.4), y2, Inches(4), Inches(0.4),
                size=15, bold=True, color=WHITE)

    chip_w = Inches(1.1)
    chip_h = Inches(0.7)
    cols   = 10
    sx     = (SLIDE_W - cols * chip_w - (cols - 1) * Inches(0.1)) / 2
    for i, d in enumerate(data):
        x = sx + i * (chip_w + Inches(0.1))
        add_rect(slide, x, y2 + Inches(0.45), chip_w, chip_h, CARD_BG)
        add_rect(slide, x, y2 + Inches(0.45), chip_w, Inches(0.06),
                 RISK_COLORS.get(d["risk_label"], WHITE))
        add_textbox(slide, d["ticker"],
                    x, y2 + Inches(0.5), chip_w, Inches(0.35),
                    size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, d["risk_label"],
                    x, y2 + Inches(0.82), chip_w, Inches(0.3),
                    size=8, color=LIGHT_GREY, align=PP_ALIGN.CENTER)


def slide_speed_breakdown(prs, data):
    """Speed label — scatter feel using positioned stat cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    accent_bar(slide)

    add_textbox(slide, "Speed Classification (CAGR-Based)",
                Inches(0.4), Inches(0.2), Inches(12), Inches(0.55),
                size=26, bold=True, color=WHITE)
    add_textbox(slide, "Based on annualized return (CAGR) from ATL to ATH",
                Inches(0.4), Inches(0.72), Inches(10), Inches(0.35),
                size=13, color=LIGHT_GREY, italic=True)

    speed_groups = {
        "Steady Climber": {"range": "CAGR < 35%/yr",  "color": GREEN},
        "Fast Mover":     {"range": "CAGR 35–70%/yr", "color": YELLOW},
        "Rocket":         {"range": "CAGR > 70%/yr",  "color": RED},
    }
    col_w   = Inches(4.0)
    col_gap = Inches(0.35)
    start_x = Inches(0.3)
    y_label = Inches(1.15)

    for col_i, (speed, meta) in enumerate(speed_groups.items()):
        x       = start_x + col_i * (col_w + col_gap)
        members = [d for d in data if d["speed_label"] == speed]

        # column header
        add_rect(slide, x, y_label, col_w, Inches(0.55), meta["color"])
        add_textbox(slide, speed,
                    x, y_label + Inches(0.05), col_w, Inches(0.45),
                    size=15, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        add_textbox(slide, meta["range"],
                    x, y_label + Inches(0.6), col_w, Inches(0.3),
                    size=10, color=LIGHT_GREY, italic=True, align=PP_ALIGN.CENTER)

        for row_i, d in enumerate(members):
            mult  = d["ath_price"] / d["atl_price"]
            import math
            cagr  = (mult ** (365.0 / d["days_between"]) - 1) * 100
            cy    = y_label + Inches(1.05) + row_i * Inches(1.18)
            add_rect(slide, x, cy, col_w, Inches(1.05), CARD_BG)
            add_rect(slide, x, cy, Inches(0.06), Inches(1.05), meta["color"])

            add_textbox(slide, d["ticker"],
                        x + Inches(0.15), cy + Inches(0.06), Inches(1.0), Inches(0.42),
                        size=16, bold=True, color=meta["color"])
            add_textbox(slide, f"CAGR: {cagr:.1f}%  ·  {mult:.1f}×  ·  {d['days_between']:,}d",
                        x + Inches(0.15), cy + Inches(0.48), col_w - Inches(0.2), Inches(0.35),
                        size=10, color=LIGHT_GREY)
            risk_color = RISK_COLORS.get(d["risk_label"], WHITE)
            add_textbox(slide, d["risk_label"],
                        x + Inches(0.15), cy + Inches(0.72), col_w - Inches(0.2), Inches(0.28),
                        size=9, color=risk_color)


def slide_ticker_detail(prs, d):
    """One slide per ticker."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)

    # left color stripe
    risk_color = RISK_COLORS.get(d["risk_label"], ACCENT)
    add_rect(slide, 0, 0, Inches(0.22), SLIDE_H, risk_color)
    accent_bar(slide, y=Inches(0))

    import math
    mult = d["ath_price"] / d["atl_price"]
    cagr = (mult ** (365.0 / d["days_between"]) - 1) * 100

    # ticker + speed
    add_textbox(slide, d["ticker"],
                Inches(0.5), Inches(0.15), Inches(4), Inches(1.0),
                size=54, bold=True, color=WHITE)
    add_textbox(slide, f"{d['speed_label']}  ·  {d['risk_label']}",
                Inches(0.5), Inches(1.05), Inches(6), Inches(0.4),
                size=16, color=risk_color)

    # 4 stat cards
    stats = [
        ("ATL Price",  f"${d['atl_price']:,.2f}",  d["atl_date"]),
        ("ATH Price",  f"${d['ath_price']:,.2f}",  d["ath_date"]),
        ("Multiplier", f"{mult:.1f}×",              f"CAGR {cagr:.1f}%/yr"),
        ("Days",       f"{d['days_between']:,}",    "ATL → ATH"),
    ]
    card_w = Inches(2.9)
    card_h = Inches(1.5)
    sx     = Inches(0.5)
    sy     = Inches(1.6)
    gap    = Inches(0.2)

    for i, (label, val, sub) in enumerate(stats):
        x = sx + i * (card_w + gap)
        add_rect(slide, x, sy, card_w, card_h, CARD_BG)
        add_rect(slide, x, sy, card_w, Inches(0.08), ACCENT)
        add_textbox(slide, label,
                    x + Inches(0.15), sy + Inches(0.1), card_w, Inches(0.35),
                    size=10, color=LIGHT_GREY)
        add_textbox(slide, val,
                    x + Inches(0.15), sy + Inches(0.42), card_w, Inches(0.65),
                    size=22, bold=True, color=WHITE)
        add_textbox(slide, sub,
                    x + Inches(0.15), sy + Inches(1.05), card_w, Inches(0.35),
                    size=10, color=LIGHT_GREY, italic=True)

    # analysis text
    add_rect(slide, Inches(0.5), Inches(3.3), Inches(12.5), Inches(0.04), ACCENT)
    add_textbox(slide, "Analysis",
                Inches(0.5), Inches(3.45), Inches(3), Inches(0.35),
                size=12, bold=True, color=ACCENT)
    add_textbox(slide, d.get("analysis", ""),
                Inches(0.5), Inches(3.85), Inches(12.4), Inches(2.8),
                size=13, color=LIGHT_GREY)


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), ACCENT)

    add_textbox(slide, "Key Takeaways",
                Inches(1), Inches(0.5), Inches(11.33), Inches(0.7),
                size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    bullets = [
        ("NVDA", "246.7× gain — most extreme multiplier in the dataset"),
        ("META & COIN", "Fastest journeys: < 1,100 days from ATL to ATH"),
        ("SPY", "Most conservative: 3.5× over 9.6 years — steady baseline"),
        ("TSLA & NVDA", "Rockets with >3,000 days still qualify — multiplier dominates speed"),
        ("Algorithm", "Current: CAGR-only. Recommended upgrade: CAGR + Max Drawdown + Sharpe (3-factor composite)"),
    ]

    for i, (ticker, text) in enumerate(bullets):
        y = Inches(1.4) + i * Inches(0.92)
        add_rect(slide, Inches(0.8), y, Inches(11.6), Inches(0.78), CARD_BG)
        add_rect(slide, Inches(0.8), y, Inches(0.08), Inches(0.78), ACCENT)
        add_textbox(slide, ticker,
                    Inches(1.1), y + Inches(0.1), Inches(1.6), Inches(0.55),
                    size=13, bold=True, color=ACCENT)
        add_textbox(slide, text,
                    Inches(2.8), y + Inches(0.12), Inches(9.4), Inches(0.55),
                    size=13, color=LIGHT_GREY)

    add_textbox(slide, "Source: Yahoo Finance  ·  Not investment advice  ·  Data as of report date",
                Inches(0), Inches(7.1), SLIDE_W, Inches(0.35),
                size=10, color=LIGHT_GREY, italic=True, align=PP_ALIGN.CENTER)


# ── main ─────────────────────────────────────────────────────────────────────

def build_deck(data: list[dict]) -> None:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs, data)
    slide_overview_table(prs, data)
    slide_multiplier_chart(prs, data)
    slide_days_chart(prs, data)
    slide_risk_breakdown(prs, data)
    slide_speed_breakdown(prs, data)

    for d in data:
        slide_ticker_detail(prs, d)

    slide_closing(prs)

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    with open(INPUT) as f:
        data = json.load(f)
    build_deck(data)
